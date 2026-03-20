from __future__ import annotations

import subprocess
import tempfile
from pathlib import Path

from docx import Document as DocxDocument
from pypdf import PdfReader
from pptx import Presentation

from sap_navigator.models import LoadedDocument, SkippedFile


SUPPORTED_SUFFIXES = {
    ".pdf",
    ".docx",
    ".pptx",
    ".txt",
    ".md",
    ".markdown",
    ".doc",
}


def load_documents(root_dir: Path) -> tuple[list[LoadedDocument], list[SkippedFile]]:
    documents: list[LoadedDocument] = []
    skipped: list[SkippedFile] = []

    if not root_dir.exists():
        return documents, [SkippedFile(path=root_dir, reason="Knowledge directory does not exist.")]

    for path in sorted(root_dir.rglob("*")):
        if not path.is_file():
            continue
        if path.name.startswith("."):
            continue

        suffix = path.suffix.lower()
        if suffix not in SUPPORTED_SUFFIXES:
            if suffix in {".one", ".onepkg"}:
                skipped.append(
                    SkippedFile(
                        path=path,
                        reason="Direct OneNote parsing is not enabled. Export the notebook/page to PDF, DOCX, or Markdown first.",
                    )
                )
            continue

        try:
            content = _load_file(path)
        except Exception as exc:  # noqa: BLE001
            skipped.append(SkippedFile(path=path, reason=str(exc)))
            continue

        normalized = _normalize_text(content)
        if not normalized.strip():
            skipped.append(SkippedFile(path=path, reason="No extractable text found."))
            continue

        documents.append(
            LoadedDocument(
                source_path=path,
                content=normalized,
                source_type=suffix.lstrip("."),
                title=path.stem,
                metadata={
                    "filename": path.name,
                    "suffix": suffix,
                    "source_path": str(path),
                },
            )
        )

    return documents, skipped


def _load_file(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return _load_pdf(path)
    if suffix == ".docx":
        return _load_docx(path)
    if suffix == ".pptx":
        return _load_pptx(path)
    if suffix in {".txt", ".md", ".markdown"}:
        return path.read_text(encoding="utf-8", errors="ignore")
    if suffix == ".doc":
        return _load_legacy_doc(path)
    raise ValueError(f"Unsupported file type: {suffix}")


def _load_pdf(path: Path) -> str:
    reader = PdfReader(str(path))
    pages: list[str] = []
    for index, page in enumerate(reader.pages, start=1):
        page_text = page.extract_text() or ""
        pages.append(f"[Page {index}]\n{page_text.strip()}")
    return "\n\n".join(pages)


def _load_docx(path: Path) -> str:
    document = DocxDocument(str(path))
    paragraphs = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
    return "\n\n".join(paragraphs)


def _load_pptx(path: Path) -> str:
    presentation = Presentation(str(path))
    slides: list[str] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = [f"[Slide {slide_index}]"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text.strip())
        slide_text = "\n".join(part for part in parts if part)
        if slide_text.strip():
            slides.append(slide_text)
    return "\n\n".join(slides)


def _load_legacy_doc(path: Path) -> str:
    textutil_result = _convert_with_textutil(path)
    if textutil_result:
        return textutil_result

    soffice_result = _convert_with_soffice(path)
    if soffice_result:
        return soffice_result

    raise ValueError("Could not convert .doc file. Install LibreOffice or use macOS textutil, or convert to DOCX/PDF first.")


def _convert_with_textutil(path: Path) -> str | None:
    command = ["textutil", "-convert", "txt", "-stdout", str(path)]
    result = subprocess.run(command, capture_output=True, text=True, check=False)
    if result.returncode == 0 and result.stdout.strip():
        return result.stdout
    return None


def _convert_with_soffice(path: Path) -> str | None:
    with tempfile.TemporaryDirectory() as temp_dir:
        out_dir = Path(temp_dir)
        command = [
            "soffice",
            "--headless",
            "--convert-to",
            "txt:Text",
            "--outdir",
            str(out_dir),
            str(path),
        ]
        result = subprocess.run(command, capture_output=True, text=True, check=False)
        if result.returncode != 0:
            return None

        converted = out_dir / f"{path.stem}.txt"
        if converted.exists():
            return converted.read_text(encoding="utf-8", errors="ignore")
    return None


def _normalize_text(text: str) -> str:
    lines = [line.rstrip() for line in text.replace("\r\n", "\n").replace("\r", "\n").split("\n")]
    compacted: list[str] = []
    blank_run = 0

    for line in lines:
        if line.strip():
            blank_run = 0
            compacted.append(line)
            continue

        blank_run += 1
        if blank_run <= 1:
            compacted.append("")

    return "\n".join(compacted).strip()

